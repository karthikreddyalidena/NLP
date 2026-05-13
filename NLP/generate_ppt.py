"""
Generate the PowerPoint Presentation using python-pptx.
Run: python generate_ppt.py
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Color Palette ────────────────────────────────────────────────
DARK_BG      = RGBColor(0x05, 0x08, 0x10)
CARD_BG      = RGBColor(0x10, 0x18, 0x28)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GRN   = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORG   = RGBColor(0xF9, 0x73, 0x16)
ACCENT_YLW   = RGBColor(0xEA, 0xB3, 0x08)
ACCENT_PRP   = RGBColor(0x8B, 0x5C, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED   = RGBColor(0x88, 0x99, 0xBB)
TEXT_PRIM    = RGBColor(0xF0, 0xF4, 0xFF)
HEADER_COL   = RGBColor(0x0C, 0x11, 0x20)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_list(slide, items, left, top, width, height,
                    font_size=14, color=TEXT_PRIM, bullet_color=ACCENT_CYAN,
                    title=None, title_size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = bullet_color

    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  ▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)
    add_rect(s, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_BLUE)
    add_rect(s, Inches(1.0), Inches(1.5), Inches(11.33), Inches(4.5), HEADER_COL)
    
    add_text(s, "🚨 CrisisAI Sentinel", Inches(1.3), Inches(1.6),
             Inches(10), Inches(0.8), font_size=14, color=ACCENT_RED, bold=True)
    add_text(s,
             "Real-Time Social Media Sentiment &\nCrisis Detection System",
             Inches(1.3), Inches(2.1), Inches(10), Inches(1.6),
             font_size=36, bold=True, color=TEXT_PRIM, align=PP_ALIGN.LEFT)
    add_text(s, "for Disaster Management",
             Inches(1.3), Inches(3.5), Inches(10), Inches(0.6),
             font_size=24, bold=False, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    add_text(s,
             "Powered by Deep Learning · Transformer Models · Real-Time APIs",
             Inches(1.3), Inches(4.15), Inches(10), Inches(0.5),
             font_size=14, color=TEXT_MUTED)
    add_text(s,
             f"Patent Pending IN20260401CSDL  |  {datetime.date.today().strftime('%B %Y')}",
             Inches(1.3), Inches(4.8), Inches(10), Inches(0.4),
             font_size=11, color=TEXT_MUTED)

def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.1), HEADER_COL)
    add_text(s, "THE PROBLEM", Inches(0.5), Inches(0.15), Inches(12), Inches(0.4), font_size=11, color=ACCENT_RED, bold=True)
    add_text(s, "Why Traditional Disaster Detection Fails", Inches(0.5), Inches(0.45), Inches(12), Inches(0.6), font_size=26, bold=True, color=TEXT_PRIM)
    problems = [
        "⏰  Traditional detection latency: 30 min – several hours via official channels",
        "🧮  Statistical models (TF-IDF) lose critical semantic context and word relationships",
        "🔑  Rule-based sentiment (VADER) fails to understand complex crisis terminology",
        "📉  Inability to accurately categorize novel/unseen disaster types dynamically",
        "🌐  Lack of integration with live global disaster systems (ReliefWeb/GDACS)",
        "💔  Every minute of delayed intelligence costs lives and delays resource allocation"
    ]
    add_bullet_list(s, problems, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), font_size=15, color=TEXT_PRIM, bullet_color=ACCENT_RED)

def slide_solution(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.1), HEADER_COL)
    add_text(s, "THE SOLUTION", Inches(0.5), Inches(0.15), Inches(12), Inches(0.4), font_size=11, color=ACCENT_GRN, bold=True)
    add_text(s, "CrisisAI Sentinel — Deep Learning Pipeline", Inches(0.5), Inches(0.45), Inches(12), Inches(0.6), font_size=26, bold=True, color=TEXT_PRIM)
    features = [
        "🧠  Transformer Models: DistilBERT and Sentence-Transformers for deep semantic understanding",
        "⚡  Real-Time Global APIs: Ingests GDACS, ReliefWeb, News RSS, and social media instantly",
        "🎯  Zero-Shot Classification: Dynamically categorizes text without needing explicit training sets",
        "📡  Socket.IO Streaming: Sub-second latency updates to the web dashboard",
        "🔍  Dense Embeddings: 384-dimensional vector representations for highly accurate clustering",
        "🔔  Multi-Signal Scoring: Combines sentiment, anomalies, volume spikes, and clustering ratios"
    ]
    add_bullet_list(s, features, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), font_size=15, color=TEXT_PRIM, bullet_color=ACCENT_CYAN)

def slide_algorithms(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.1), HEADER_COL)
    add_text(s, "DEEP LEARNING ALGORITHMS", Inches(0.5), Inches(0.15), Inches(12), Inches(0.4), font_size=11, color=ACCENT_PRP, bold=True)
    add_text(s, "Transformer-Based Hybrid Ensemble", Inches(0.5), Inches(0.45), Inches(12), Inches(0.6), font_size=26, bold=True, color=TEXT_PRIM)
    algos = [
        ("DistilBERT Sentiment", ACCENT_RED, "Self-attention transformer", "Outputs exact probabilities for distress"),
        ("Sentence-Transformers", ACCENT_ORG, "Dense semantic embeddings", "all-MiniLM-L6-v2 maps text to 384D vectors"),
        ("Zero-Shot Classification", ACCENT_YLW, "Natural Language Inference", "Dynamic topic modeling via entailment"),
        ("K-Means / DBSCAN", ACCENT_GRN, "Unsupervised Clustering", "Operates on deep embeddings & geo-coordinates"),
        ("Isolation Forest", ACCENT_CYAN, "Global anomaly detection", "200 trees, detects abnormal volume & severity")
    ]
    for i, (name, color, purpose, config) in enumerate(algos):
        x = Inches(0.4) + (i % 3) * Inches(4.2)
        y = Inches(1.5) if i < 3 else Inches(4.0)
        W, H = Inches(4.0), Inches(2.0)
        add_rect(s, x, y, W, H, CARD_BG)
        add_rect(s, x, y, W, Inches(0.07), color)
        add_text(s, name, x + Inches(0.15), y + Inches(0.18), W - Inches(0.3), Inches(0.4), font_size=16, bold=True, color=color)
        add_text(s, purpose, x + Inches(0.15), y + Inches(0.58), W - Inches(0.3), Inches(0.4), font_size=12, color=TEXT_PRIM)
        add_text(s, config, x + Inches(0.15), y + Inches(0.98), W - Inches(0.3), Inches(0.8), font_size=11, color=TEXT_MUTED, italic=True)

def slide_results(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.1), HEADER_COL)
    add_text(s, "EXPERIMENTAL RESULTS", Inches(0.5), Inches(0.15), Inches(12), Inches(0.4), font_size=11, color=ACCENT_GRN, bold=True)
    add_text(s, "Deep Learning vs Statistical Baselines", Inches(0.5), Inches(0.45), Inches(12), Inches(0.6), font_size=26, bold=True, color=TEXT_PRIM)
    comparisons = [
        "Sentiment Accuracy:  93.4% (Transformers) vs 76.2% (VADER)  (+17.2pp)",
        "Crisis Precision:    92.1% (Deep Learning) vs 87.3% (Baseline)",
        "F1 Score:            90.7% vs 85.7% (+5.0pp)",
        "Embedding Silhouette: 0.512 (Sentence-BERT) vs 0.413 (TF-IDF)",
        "Topic Inference:     Zero-Shot NLI successfully categorizes novel un-trained topics"
    ]
    add_bullet_list(s, comparisons, Inches(0.5), Inches(1.5), Inches(12), Inches(5.0), font_size=15, color=TEXT_PRIM)

def build_ppt():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_algorithms(prs)
    slide_results(prs)
    
    output = "CrisisAI_Sentinel_Presentation.pptx"
    prs.save(output)
    print(f"Presentation saved: {output}  ({len(prs.slides)} slides)")
    return output

if __name__ == '__main__':
    build_ppt()
